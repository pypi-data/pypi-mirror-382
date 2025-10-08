from __future__ import annotations
from pathlib import Path
import io

from bs4 import BeautifulSoup  # html
from pypdf import PdfReader    # pdf
from docx import Document      # docx
from pptx import Presentation  # pptx
from striprtf.striprtf import rtf_to_text  # rtf

TEXT_EXTS = {".txt", ".md", ".markdown", ".log", ".json", ".yaml", ".yml", ".csv", ".tsv"}
HTML_EXTS = {".html", ".htm"}

def _decode_utf8(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")

def extract_text_from_bytes(path: Path, data: bytes) -> str:
    ext = path.suffix.lower()

    # Obvious text-like files
    if ext in TEXT_EXTS:
        return _decode_utf8(data)

    if ext == ".rtf":
        return rtf_to_text(_decode_utf8(data))

    if ext in HTML_EXTS:
        html = _decode_utf8(data)
        soup = BeautifulSoup(html, "html5lib")
        return soup.get_text(separator="\n")

    if ext == ".pdf":
        reader = PdfReader(io.BytesIO(data))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n".join(parts).strip()

    if ext == ".docx":
        doc = Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs]
        # tables text
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join([p for p in parts if p]).strip()

    if ext == ".pptx":
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(parts).strip()

    # Fallback: best-effort UTF-8 decode
    return _decode_utf8(data)
